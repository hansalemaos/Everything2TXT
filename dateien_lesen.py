# -*- coding: utf-8 -*-
import docx2txt
from pdfminer.high_level import extract_text
from pptx import Presentation
from bs4 import BeautifulSoup
from epubextract import epub2txt
from xlsx2html import xlsx2html
import tempfile
from tkinter import Tk
from tkinter.filedialog import askopenfilename
import re


def create_temp_file(ending):
    fp = tempfile.TemporaryFile(suffix=ending, delete=False)
    return fp.name


def powerpointlesen(pfad):
    prs = Presentation(pfad)
    ganzertext = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text"):
                    ganzertext = ganzertext + "\n" + shape.text
            except Exception as Fehler:
                print(Fehler)
    return ganzertext


def docxlesen(pfad):
    return docx2txt.process(pfad)


def txtdateien_lesen(pfad):
    try:
        with open(pfad, mode="rb") as f:
            dateiohnehtml = f.read()
        dateiohnehtml = (
            b"""<!DOCTYPE html><html><body><p>"""
            + dateiohnehtml
            + b"""</p></body></html>"""
        )
        soup = BeautifulSoup(dateiohnehtml, "lxml")
        soup = soup.text
        return soup.strip()
    except Exception as Fehler:
        print(Fehler)
        with open(pfad, mode="r", encoding="utf-8") as f:
            dateiohnehtml = f.read()
        return dateiohnehtml


def html_htm_dateien_lesen(pfad):
    try:
        with open(pfad, mode="rb") as f:
            dateiohnehtml = f.read()
        soup = BeautifulSoup(dateiohnehtml, "lxml")
        soup = soup.text
        soup = soup.strip()
        return soup
    except Exception as Fehler:
        print(Fehler)


def pdf_datei_lesen(pfad):
    return extract_text(pfad)


def xlsx_datei_einlesen(pfad):
    tmpdatei = create_temp_file(ending="html")
    xlsx2html(pfad, tmpdatei)
    text = html_htm_dateien_lesen(tmpdatei)
    return text


def dateienauslesen(pfad):
    if str(pfad).endswith("pptx"):
        text = powerpointlesen(pfad)
        return text
    elif str(pfad).endswith("docx"):
        text = docxlesen(pfad)
        return text
    elif str(pfad).endswith("html") or str(pfad).endswith("htm"):
        text = txtdateien_lesen(pfad)
        return text

    elif str(pfad).endswith("pdf"):
        text = pdf_datei_lesen(pfad)
        return text

    elif str(pfad).endswith("epub"):
        text = epub2txt(pfad)
        text = text.convert()
        return text

    elif str(pfad).endswith("xlsx"):
        text = xlsx_datei_einlesen(pfad)
        return text
    else:
        text = txtdateien_lesen(pfad)
        return text


def datei_auswaehlen_mit_tkinter():
    Tk().withdraw()
    dateiname = askopenfilename()
    ausgabeordner = re.sub(r"/[^/]+\.\w+$", "", dateiname)
    ausgabedatei = re.sub(r"^.*(/[^/]+)\.\w{,8}", "\g<1>.txt", dateiname)
    ausgabedatei = ausgabeordner + ausgabedatei
    return dateiname, ausgabedatei


if __name__ == "__main__":
    dateiname, ausgabedatei = datei_auswaehlen_mit_tkinter()
    textzumspeichern = dateienauslesen(dateiname)
    if not str(dateiname).endswith(".txt"):
        with open(ausgabedatei, mode="w", encoding="utf-8") as f:
            if isinstance(textzumspeichern, str):
                f.write(textzumspeichern)
            if isinstance(textzumspeichern, list):
                textzumspeichern = "\n".join(textzumspeichern)
                f.write(textzumspeichern)

    print(textzumspeichern)
